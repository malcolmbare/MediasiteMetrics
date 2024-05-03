import os
import json
from datetime import datetime
import requests
import zipfile
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from ReqAndAuth import mediasite
from ReqAndAuth import panopto
import xml.etree.ElementTree as ET
import codecs
import boto3
import cv2
import time
import copy


class mediasitePublish2Go:

    def createPublishToGo(self):
        #Create Endpoints
        addPublishToGo = ms.serverURL + "/Presentations('{0}')/AddPublishToGo".format(self.id)
        submitPublishToGo = ms.serverURL + "/Presentations('{0}')/SubmitPublishToGo".format(self.id)
        #Some videos already have PublishToGo Content available. This means that the firs
        self.checkStatus(addPublishToGo, "AddPublishToGo")
        self.checkStatus(submitPublishToGo, "SubmitPublishToGo")
    def downloadPublishToGo(self):
        #Create Endpoint
        getPublishToGo = ms.serverURL + "/Presentations('{0}')/PublishToGoContent".format(self.id)
        #Get the Download URL
        downloadURL = self.awaitDownloadURL(getPublishToGo)
        self.downloadURL = downloadURL
    def awaitDownloadURL(self, endpoint):
        response = requests.get(endpoint, headers=ms.header)
        result = response.json()["DownloadUrl"]
        while result == None:
            print("Checking...")
            time.sleep(15)
            result = response.json()["DownloadUrl"]
        print("Download URL: {0}".format(result.json()["DownloadUrl"]))
        return result
    def checkStatus(self, endpoint, name):
        if self.session.post(endpoint, headers=ms.header).status_code <= 300:
            print("{0} Successful".format(name))
    def __init__(self, id):
        self.id = id
        self.session = requests.Session()
        self.createPublishToGo()
        self.download = self.session.get(self.downloadURL, headers=ms.header)
        self.zipfile = zipfile.ZipFile(BytesIO(self.download.content))
        self.path = "MSFiles/{0}".format(self.id)
        self.zipfile.extractall(self.path)
        print("Complete")

class mediasitePresentation:
    def getSlideNumbers(self, path):
        tree = ET.parse(path + "/" + self.xml_file_names[0])
        root = tree.getroot()
        return [i.text for i in root.iter("Number")]
    def getSlideTimes(self, path):
        tree = ET.parse(path + "/" + self.xml_file_names[0])
        root = tree.getroot()
        rawTimes = [i.text for i in root.iter("Time")]
        result = [int(i) / 1000 for i in rawTimes]
        return result
    def makeSlideTuples(self):
        if len(self.image_stream) == len(self.xml_slide_times):
            return [
                (self.image_stream[i], self.xml_slide_times[i])
                for i in range(len(self.xml_slide_times))
            ]
    def makeSlideshow(self):
        template = cv2.imread(self.image_stream[0])
        h, w, _ = template.shape
        prs = Presentation()
        prs.slide_width = w * 10000
        prs.slide_height = h * 10000
        blank_slide_layout = prs.slide_layouts[6]
        for file, duration in self.slideTuples:
            slide = prs.slides.add_slide(blank_slide_layout)
            pic = slide.shapes.add_picture(file,0,0, width=prs.slide_width)
        prs.save(self.nextPath+"/"+'output.pptx')
    def createManifest(self, file_path, manifest_file_name, newTitle, description):
        file_name = os.path.basename(file_path)
        with open(self.manifestFileTemplate) as fr:
            template = fr.read()
            content = (
                template.replace("{Title}", newTitle)
                .replace(
                    "{Description}",
                    description,
                )
                .replace("{Filename}", file_name)
                .replace(
                    "{Date}", datetime.utcnow().strftime("%Y-%m-%dT%H:%M:%S.%f-00:00")
                )
            )
            print(content)
        with codecs.open(manifest_file_name, "w", "utf-8") as fw:
            fw.write(content)
    def addSlidesToManifest(self, manifest_file_name):
        tree = ET.parse(manifest_file_name)
        ET.register_namespace("", "http://tempuri.org/UniversalCaptureSpecification/v1")
        root= tree.getroot()
        Presentations = ET.Element("Presentations")
        presentation= ET.SubElement(Presentations, 'Presentation')
        presentation= self.addStart(presentation)
        presentation = self.addFile(presentation)
        presentation = self.addSlideChanges(presentation)
        root.append(Presentations)
        ET.indent(tree, '  ')
        tree.write(manifest_file_name)
    def addStart(self, presentation):
        start = ET.SubElement(presentation, 'Start')
        start.text = 'PT0S'
        return presentation
    def addFile(self, presentation):
        file = ET.SubElement(presentation, 'File')
        file.text = 'output.pptx'
        return presentation
    def addSlideChanges(self,presentation):
        slideChanges = ET.SubElement(presentation, 'SlideChanges')
        for i in range(len(self.slideTuples)):
            slideChange = ET.SubElement(slideChanges, 'SlideChange')
            slideNumber = ET.SubElement(slideChange, 'SlideNumber')
            slideNumber.text = str(i +1)
            time = ET.SubElement(slideChange, 'Time')
            timetxt = str(self.slideTuples[i][1])
            if "." in timetxt:
                timetxt=timetxt.split(".")[0]
            time.text = "PT" + timetxt + "S"
        return presentation
    def checkRound(self, frames2Insert):
        if self.roundedLast == False:
            self.roundedLast = True
            frameInsert = int(round(frames2Insert))
        else:
            self.roundedLast = False
            frameInsert = int(frames2Insert)
        return frameInsert
    def __init__(self, path, newTitle, description):
        self.nextPath = path + "/content"
        self.file_names = [
            file for file in os.listdir(self.nextPath) if file.endswith(".wmv")
        ]
        self.file_names += [
            file
            for file in os.listdir(self.nextPath)
            if file.endswith(".mp4")
            if "output" not in file
        ]
        self.file_path = self.nextPath + "/" + self.file_names[0]
        self.image_stream = sorted(
            [
                self.nextPath + "/" + file
                for file in os.listdir(self.nextPath)
                if file.endswith("full.jpg")
            ]
        )
        self.xml_file_names = [
            file for file in os.listdir(path) if file.endswith(".xml")
        ]
        self.xml_slide_numbers = self.getSlideNumbers(path)
        self.xml_slide_times = self.getSlideTimes(path)
        self.slideTuples = self.makeSlideTuples()
        self.roundedLast = False
        self.slideStream_path = self.nextPath + "/output.pptx"
        self.manifestFileTemplate = "upload_manifest_template.xml"
        self.manifestFileName = "upload_manifest_generated.xml"
        self.createManifest(self.file_path, self.manifestFileName, newTitle, description)
        if len(self.slideTuples) > 0:
            self.addSlidesToManifest(self.manifestFileName)
            self.makeSlideshow()

class makePanoptoFolder:
    def checkIfExists(self, name):
        resp = requests.get(self.createFolderEndpoint +"/search?searchQuery={0}".format(name), headers=pan.header)
        if len(resp.json()["Results"])>0:
            self.newFolderId = resp.json()["Results"][0]["Id"]
            print("Folder Exists")
        else:
            self.payload = json.dumps(
            {"Name": name, "Description": "Folder Working?", "Parent": self.mainFolder})
            self.createFolder = requests.post(
                self.createFolderEndpoint, headers=pan.header, data=self.payload
            )
            self.newFolderId = self.createFolder.json()["Id"]
            print("Folder Created")
    def __init__(self, name, mainFolder):
        self.mainFolder = mainFolder
        self.createFolderEndpoint = pan.serverURL + "/Panopto/api/v1/folders"
        self.checkIfExists(name)

class uploadPanopto:
    def createSession(self):
        payload = json.dumps({"FolderId": self.folderId})
        self.uploadSessionResp = requests.post(
            pan.serverURL + "/Panopto/PublicAPI/Rest/sessionUpload",
            headers=pan.header,
            data=payload,
        )
        self.uploadID = self.uploadSessionResp.json()["ID"]
        self.uploadTarget = self.uploadSessionResp.json()["UploadTarget"]
    def multipartUpload(self, uploadTarget, file_path):
        elements = uploadTarget.split("/")
        service_endpoint = "/".join(elements[0:-2:])
        bucket = elements[-2]
        prefix = elements[-1]
        object_key = "{0}/{1}".format(prefix, os.path.basename(file_path))
        s3 = boto3.session.Session().client(
            service_name="s3",
            endpoint_url=service_endpoint,
            verify=False,
            aws_access_key_id="dummy",
            aws_secret_access_key="dummy",
        )
        mpu = s3.create_multipart_upload(Bucket=bucket, Key=object_key)
        mpu_id = mpu["UploadId"]
        parts = []
        uploaded_bytes = 0
        total_bytes = os.stat(file_path).st_size
        with open(file_path, "rb") as f:
            i = 1
            while True:
                data = f.read(self.PART_SIZE)
                if not len(data):
                    break
                part = s3.upload_part(
                    Body=data,
                    Bucket=bucket,
                    Key=object_key,
                    UploadId=mpu_id,
                    PartNumber=i
                )
                parts.append({"PartNumber": i, "ETag": part["ETag"]})
                uploaded_bytes += len(data)
                print(
                    "  -- {0} of {1} bytes uploaded".format(uploaded_bytes, total_bytes)
                )
                i += 1
        result = s3.complete_multipart_upload(
            Bucket=bucket,
            Key=object_key,
            UploadId=mpu_id,
            MultipartUpload={"Parts": parts},
        )
    def finishUpload(self, uploadSessionResp):
        upload_id = self.uploadID
        while True:
            print(
                "Calling PUT PublicAPI/REST/sessionUpload/{0} endpoint".format(
                    upload_id
                )
            )
            url = "{0}/Panopto/PublicAPI/REST/sessionUpload/{1}".format(
                pan.serverURL, upload_id
            )
            payload = copy.copy(uploadSessionResp.json())
            payload["State"] = 1  # Upload Completed
            resp = requests.put(url, json=payload, headers=pan.header)
            if resp.status_code // 100 == 2:
                print("Upload completed")
                break
    def __init__(self, folderId, file_path, slideStream_path):
        self.PART_SIZE = 5 * 1024 * 1024
        self.manifestFileTemplate = "upload_manifest_template.xml"
        self.manifestFileName = "upload_manifest_generated.xml"
        self.folderId = folderId
        self.createSession()
        self.multipartUpload(self.uploadTarget, file_path)
        if os.path.exists(slideStream_path):
            self.multipartUpload(self.uploadTarget, slideStream_path)
        self.multipartUpload(self.uploadTarget, "upload_manifest_generated.xml")
        self.finishUpload(self.uploadSessionResp)

class uploadFile:
    def __init__ (self, row):
        self.id = row[-3]
        self.folder = row[1]
        self.subfolder = row[-1]
        self.newtitle = row[-2]
        self.description = "The following video was originally titled {0}.".format(row[0])
        self.description += " We believe this was a {0} that was recorded on {1} at {2}.".format(row[2], row[4], row[3])
        self.description += " The original recording had {0} total views, {1} unique views, and {2} views since 2020.".format(row[5], row[6], row[7])

class prog:
    def __init__(self, id, folder, subfolder, newtitle, description):
        msp2g = mediasitePublish2Go(id)
        msx = mediasitePresentation(msp2g.path, newtitle, description)
        pan = panopto()
        pf = makePanoptoFolder(folder,"1d52671a-f0f2-4873-b088-b163014744c6")
        pfs = makePanoptoFolder(subfolder, pf.newFolderId)
        up = uploadPanopto(pfs.newFolderId, msx.file_path, msx.slideStream_path)

rows = []
with open('admissionsTest.csv', 'r') as file:
    reader = csv.reader(file)
    next(reader)
    for row in reader:
        if row[-1] != "":
            rows.append(uploadFile(row))
